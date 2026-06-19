"""
文件解析模块
支持 PDF、Word、TXT、Markdown 等格式的文档解析
"""
from pathlib import Path
from typing import Optional, List, Dict
from utils.logger import get_logger

logger = get_logger("file_parser")


class FileParser:
    """文件解析器，将不同格式文档统一转为纯文本"""
    
    SUPPORTED_FORMATS = {".pdf", ".docx", ".doc", ".pptx", ".ppt", ".xlsx", ".xls", ".txt", ".md", ".csv"}
    
    @staticmethod
    def parse(file_path: str) -> str:
        """
        根据文件扩展名自动选择解析方式
        
        Args:
            file_path: 文件路径
        
        Returns:
            解析后的纯文本内容
        """
        path = Path(file_path)
        if not path.exists():
            raise FileNotFoundError(f"文件不存在: {file_path}")
        
        suffix = path.suffix.lower()
        logger.info(f"解析文件: {file_path}, 类型: {suffix}")
        
        parsers = {
            ".pdf": FileParser._parse_pdf,
            ".docx": FileParser._parse_docx,
            ".doc": FileParser._parse_docx,
            ".pptx": FileParser._parse_pptx,
            ".ppt": FileParser._parse_pptx,
            ".xlsx": FileParser._parse_excel,
            ".xls": FileParser._parse_excel,
            ".txt": FileParser._parse_text,
            ".md": FileParser._parse_text,
            ".csv": FileParser._parse_csv,
        }
        
        parser = parsers.get(suffix)
        if parser is None:
            raise ValueError(f"不支持的文件格式: {suffix}")
        
        return parser(file_path)
    
    @staticmethod
    def parse_with_pages(file_path: str) -> List[Dict]:
        """
        解析文件，返回带页码信息的文本片段列表
        
        Args:
            file_path: 文件路径
        
        Returns:
            [{"text": "页面文本", "page": 1}, {"text": "全文", "page": None}]
            PDF 返回逐页结果，其他格式 page 为 None
        """
        path = Path(file_path)
        if not path.exists():
            raise FileNotFoundError(f"文件不存在: {file_path}")
        
        suffix = path.suffix.lower()
        
        if suffix == ".pdf":
            return FileParser._parse_pdf_with_pages(file_path)
        if suffix in (".pptx", ".ppt"):
            return FileParser._parse_pptx_with_pages(file_path)
        if suffix in (".xlsx", ".xls"):
            return FileParser._parse_excel_with_pages(file_path)
        
        # 其他格式：单块文本，无页码
        text = FileParser.parse(file_path)
        return [{"text": text, "page": None}]
    
    @staticmethod
    def _parse_pdf_with_pages(file_path: str) -> List[Dict]:
        """解析 PDF 文件，返回逐页文本"""
        try:
            import fitz  # PyMuPDF
            doc = fitz.open(file_path)
            pages = []
            for i, page in enumerate(doc):
                text = page.get_text()
                if text.strip():
                    pages.append({"text": text, "page": i + 1})
            doc.close()
            logger.info(f"PDF 解析完成: {len(pages)} 页有文本内容")
            return pages
        except ImportError:
            raise ImportError("请安装 PyMuPDF: pip install PyMuPDF")
    
    @staticmethod
    def _parse_pdf(file_path: str) -> str:
        """解析 PDF 文件（纯文本，不保留页码）"""
        try:
            import fitz  # PyMuPDF
            doc = fitz.open(file_path)
            text_parts = []
            for page in doc:
                text_parts.append(page.get_text())
            doc.close()
            return "\n".join(text_parts)
        except ImportError:
            raise ImportError("请安装 PyMuPDF: pip install PyMuPDF")
    
    @staticmethod
    def _parse_docx(file_path: str) -> str:
        """解析 Word 文档"""
        try:
            from docx import Document
            doc = Document(file_path)
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            return "\n".join(paragraphs)
        except ImportError:
            raise ImportError("请安装 python-docx: pip install python-docx")
    
    @staticmethod
    def _parse_text(file_path: str) -> str:
        """解析纯文本和 Markdown 文件（自动检测编码）"""
        import chardet
        with open(file_path, "rb") as f:
            raw = f.read()
        detected = chardet.detect(raw)
        encoding = detected.get("encoding") or "utf-8"
        # chardet 有时返回 "ascii"，实际可能是 utf-8
        if encoding.lower() == "ascii":
            encoding = "utf-8"
        logger.debug(f"文件 {file_path} 检测编码: {encoding} (置信度: {detected.get('confidence', 0):.0%})")
        return raw.decode(encoding, errors="replace")
    
    @staticmethod
    def _parse_csv(file_path: str) -> str:
        """解析 CSV 文件"""
        try:
            import pandas as pd
            df = pd.read_csv(file_path)
            return df.to_string(index=False)
        except ImportError:
            raise ImportError("请安装 pandas: pip install pandas")

    @staticmethod
    def _df_to_markdown(df) -> str:
        """将 DataFrame 转为 Markdown 表格格式"""
        if df.empty:
            return ""
        headers = [str(h).strip() for h in df.columns.tolist()]
        rows = []
        for _, row in df.iterrows():
            rows.append([str(v).strip() if v is not None else "" for v in row.tolist()])
        # 计算每列最大宽度
        col_widths = [len(h) for h in headers]
        for row in rows:
            for i, cell in enumerate(row):
                col_widths[i] = max(col_widths[i], len(cell))
        # 构建 Markdown 表格
        header_line = "| " + " | ".join(h.ljust(col_widths[i]) for i, h in enumerate(headers)) + " |"
        separator = "| " + " | ".join("-" * max(3, col_widths[i]) for i in range(len(headers))) + " |"
        data_lines = []
        for row in rows:
            line = "| " + " | ".join(cell.ljust(col_widths[i]) for i, cell in enumerate(row)) + " |"
            data_lines.append(line)
        return "\n".join([header_line, separator] + data_lines)

    @staticmethod
    def _parse_excel(file_path: str) -> str:
        """解析 Excel 文件（纯文本，所有 sheet 合并，Markdown 表格格式）"""
        try:
            import pandas as pd
        except ImportError:
            raise ImportError("请安装 pandas: pip install pandas")

        xls = pd.ExcelFile(file_path)
        text_parts = []
        for sheet_name in xls.sheet_names:
            df = pd.read_excel(xls, sheet_name=sheet_name)
            if df.empty:
                continue
            text_parts.append(f"【{sheet_name}】")
            text_parts.append(FileParser._df_to_markdown(df))
        xls.close()
        return "\n".join(text_parts)

    @staticmethod
    def _parse_excel_with_pages(file_path: str) -> List[Dict]:
        """解析 Excel 文件，每个 sheet 作为一页（Markdown 表格格式）"""
        try:
            import pandas as pd
        except ImportError:
            raise ImportError("请安装 pandas: pip install pandas")

        xls = pd.ExcelFile(file_path)
        pages = []
        for i, sheet_name in enumerate(xls.sheet_names, 1):
            df = pd.read_excel(xls, sheet_name=sheet_name)
            if df.empty:
                continue
            text = f"【{sheet_name}】\n" + FileParser._df_to_markdown(df)
            pages.append({"text": text, "page": i})
        xls.close()
        logger.info(f"Excel 解析完成: {len(pages)} 个 sheet 有数据")
        return pages

    @staticmethod
    def _parse_pptx(file_path: str) -> str:
        """解析 PPT 文件（纯文本，不保留页码）"""
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError("请安装 python-pptx: pip install python-pptx")

        prs = Presentation(file_path)
        text_parts = []
        for slide_idx, slide in enumerate(prs.slides, 1):
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            text_parts.append(text)
                # 表格内容
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = " | ".join(
                            cell.text.strip() for cell in row.cells if cell.text.strip()
                        )
                        if row_text:
                            text_parts.append(row_text)
        return "\n".join(text_parts)

    @staticmethod
    def _parse_pptx_with_pages(file_path: str) -> List[Dict]:
        """解析 PPT 文件，返回逐页文本（每页对应一张幻灯片）"""
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError("请安装 python-pptx: pip install python-pptx")

        prs = Presentation(file_path)
        pages = []
        for slide_idx, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_texts.append(text)
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = " | ".join(
                            cell.text.strip() for cell in row.cells if cell.text.strip()
                        )
                        if row_text:
                            slide_texts.append(row_text)
            if slide_texts:
                pages.append({"text": "\n".join(slide_texts), "page": slide_idx})
        logger.info(f"PPT 解析完成: {len(pages)} 页有文本内容")
        return pages

    @staticmethod
    def extract_images_from_pdf(
        file_path: str, output_dir: str
    ) -> List[Dict]:
        """从 PDF 中提取图片并保存到文件

        Args:
            file_path: PDF 文件路径
            output_dir: 图片保存目录

        Returns:
            [{"page": 1, "path": "/path/img_p1_1.png", "b64": "iVBOR..."}]
        """
        import base64
        import os

        try:
            import fitz  # PyMuPDF
        except ImportError:
            raise ImportError("请安装 PyMuPDF: pip install PyMuPDF")

        doc = fitz.open(file_path)
        os.makedirs(output_dir, exist_ok=True)
        images = []

        for page_idx, page in enumerate(doc):
            img_list = page.get_images(full=True)
            for img_idx, img_info in enumerate(img_list):
                xref = img_info[0]
                try:
                    base_image = doc.extract_image(xref)
                    image_bytes = base_image["image"]
                    ext = base_image.get("ext", "png")

                    # 保存到文件
                    img_name = f"img_p{page_idx+1}_{img_idx+1}.{ext}"
                    img_path = os.path.join(output_dir, img_name)
                    with open(img_path, "wb") as f:
                        f.write(image_bytes)

                    images.append({
                        "page": page_idx + 1,
                        "path": img_path,
                        "b64": base64.b64encode(image_bytes).decode(),
                    })
                except Exception as e:
                    logger.warning(
                        f"提取图片失败: page={page_idx+1}, "
                        f"img={img_idx+1}, error={e}"
                    )
                    continue

        doc.close()
        logger.info(f"PDF 图片提取完成: {len(images)} 张图片")
        return images

    @staticmethod
    def extract_images_from_docx(
        file_path: str, output_dir: str
    ) -> List[Dict]:
        """从 Word 文档中提取图片并保存到文件

        Args:
            file_path: Word 文件路径
            output_dir: 图片保存目录

        Returns:
            [{"page": 0, "path": "/path/img.png", "b64": "iVBOR..."}]
        """
        import base64
        import os

        try:
            from docx import Document
        except ImportError:
            raise ImportError("请安装 python-docx: pip install python-docx")

        doc = Document(file_path)
        os.makedirs(output_dir, exist_ok=True)
        images = []

        for rel in doc.part.rels.values():
            if "image" in rel.reltype:
                try:
                    img_data = rel.target_part.blob
                    # 从 content_type 推断扩展名
                    content_type = rel.target_part.content_type
                    ext_map = {
                        "image/png": "png",
                        "image/jpeg": "jpg",
                        "image/gif": "gif",
                        "image/webp": "webp",
                        "image/bmp": "bmp",
                        "image/tiff": "tiff",
                        "image/x-emf": "emf",
                        "image/x-wmf": "wmf",
                    }
                    ext = ext_map.get(content_type, "png")

                    img_name = f"img_{rel.rId}.{ext}"
                    img_path = os.path.join(output_dir, img_name)
                    with open(img_path, "wb") as f:
                        f.write(img_data)

                    images.append({
                        "page": 0,  # Word 无精确页码
                        "path": img_path,
                        "b64": base64.b64encode(img_data).decode(),
                    })
                except Exception as e:
                    logger.warning(f"Word 图片提取失败: relId={rel.rId}, error={e}")
                    continue

        logger.info(f"Word 图片提取完成: {len(images)} 张图片")
        return images

    @staticmethod
    def extract_images_from_pptx(
        file_path: str, output_dir: str
    ) -> List[Dict]:
        """从 PPT 中提取图片并保存到文件

        Args:
            file_path: PPT 文件路径
            output_dir: 图片保存目录

        Returns:
            [{"page": 1, "path": "/path/img.png", "b64": "iVBOR..."}]
        """
        import base64
        import os

        try:
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE
        except ImportError:
            raise ImportError("请安装 python-pptx: pip install python-pptx")

        prs = Presentation(file_path)
        os.makedirs(output_dir, exist_ok=True)
        images = []

        for slide_idx, slide in enumerate(prs.slides, 1):
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        img_data = shape.image.blob
                        content_type = shape.image.content_type
                        ext_map = {
                            "image/png": "png",
                            "image/jpeg": "jpg",
                            "image/gif": "gif",
                            "image/webp": "webp",
                            "image/bmp": "bmp",
                        }
                        ext = ext_map.get(content_type, "png")

                        img_name = f"img_s{slide_idx}_{shape.shape_id}.{ext}"
                        img_path = os.path.join(output_dir, img_name)
                        with open(img_path, "wb") as f:
                            f.write(img_data)

                        images.append({
                            "page": slide_idx,
                            "path": img_path,
                            "b64": base64.b64encode(img_data).decode(),
                        })
                    except Exception as e:
                        logger.warning(
                            f"PPT 图片提取失败: slide={slide_idx}, "
                            f"shape={shape.shape_id}, error={e}"
                        )
                        continue

        logger.info(f"PPT 图片提取完成: {len(images)} 张图片")
        return images